import os
from llama_parse import LlamaParse
from pptx import Presentation
from server.logger.logger_config import my_logger as logger

USE_LLAMA_PARSE = int(os.getenv('USE_LLAMA_PARSE'))
LLAMA_CLOUD_API_KEY = os.getenv('LLAMA_CLOUD_API_KEY')


class AsyncPptxLoader:
    def __init__(self, file_path: str) -> None:
        logger.info(f"[FILE LOADER] init pptx, file_path: '{file_path}'")
        self.file_path = file_path

    async def get_content(self) -> str:
        try:
            content = ''

            if USE_LLAMA_PARSE:
                parser = LlamaParse(
                    api_key=LLAMA_CLOUD_API_KEY,
                    result_type="markdown",
                )

                text_vec = []

                import nest_asyncio
                nest_asyncio.apply()

                documents = parser.load_data(self.file_path)
                for doc in documents:
                    text_vec.append(doc.text)
                content = "\n\n".join(text_vec)
            else:
                # Load the presentation
                prs = Presentation(self.file_path)
                # Initialize a list to hold markdown parts
                markdown_parts = []

                # Process each slide in the presentation
                for slide_number, slide in enumerate(prs.slides, start=1):
                    # Add a slide header
                    markdown_parts.append(f"## Slide {slide_number}\n")

                    # Process each shape in the slide
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        text_frame = shape.text_frame

                        # Process each paragraph in the text frame
                        for paragraph in text_frame.paragraphs:
                            # Combine the runs in the paragraph to form a full text
                            text_runs = [run.text for run in paragraph.runs]
                            paragraph_text = ''.join(text_runs).strip()

                            # Convert the text into a markdown bullet point
                            if paragraph_text:
                                markdown_parts.append(f"- {paragraph_text}\n")

                if markdown_parts:
                    # Join all parts to form the final markdown text
                    content = ''.join(markdown_parts)

            if not content:
                logger.warning(f"file_path: '{self.file_path}' is empty!")
            return content
        except Exception as e:
            logger.error(f"get_content is failed, exception: {e}")
            return ''
